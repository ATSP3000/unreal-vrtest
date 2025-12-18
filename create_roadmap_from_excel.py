#!/usr/bin/env python3
"""
Generates a PowerPoint roadmap from an Excel data file.

Usage:
    python create_roadmap_from_excel.py <excel_file> [output_file]

Example:
    python create_roadmap_from_excel.py roadmap_data.xlsx my_roadmap.pptx
"""

import sys
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# Slide dimensions (standard 16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def hex_to_rgb(hex_color):
    """Convert hex color string to RGBColor."""
    hex_color = hex_color.lstrip('#')
    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16)
    )


def add_text_box(slide, left, top, width, height, text, font_size=10, bold=False, 
                 font_color=None, fill_color=None, alignment=PP_ALIGN.LEFT, italic=False):
    """Add a text box to the slide."""
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    
    p = tf.paragraphs[0]
    p.text = str(text)
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.italic = italic
    p.alignment = alignment
    
    if font_color:
        p.font.color.rgb = font_color
    
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    
    return shape


def add_rectangle(slide, left, top, width, height, fill_color=None, line_color=None):
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    
    return shape


def add_triangle_milestone(slide, left, top, size=Inches(0.12), color=None):
    """Add a triangle milestone marker."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def load_excel_data(excel_path):
    """Load all data from the Excel file."""
    data = {}
    
    # Load Settings
    settings_df = pd.read_excel(excel_path, sheet_name="Settings")
    data['settings'] = dict(zip(settings_df['Setting'], settings_df['Value']))
    
    # Load Timeline
    data['timeline'] = pd.read_excel(excel_path, sheet_name="Timeline")
    
    # Load Goals
    data['goals'] = pd.read_excel(excel_path, sheet_name="Goals")
    
    # Load Rows
    data['rows'] = pd.read_excel(excel_path, sheet_name="Rows")
    
    # Load Milestones
    data['milestones'] = pd.read_excel(excel_path, sheet_name="Milestones")
    
    # Load Use Cases
    data['usecases'] = pd.read_excel(excel_path, sheet_name="UseCases")
    
    return data


def create_roadmap_from_data(data):
    """Create the roadmap presentation from loaded data."""
    settings = data['settings']
    timeline_df = data['timeline']
    goals_df = data['goals']
    rows_df = data['rows']
    milestones_df = data['milestones']
    usecases_df = data['usecases']
    
    # Parse colors from settings
    colors = {
        'navy': hex_to_rgb(settings.get('Navy Color', '1A1A4E')),
        'pink': hex_to_rgb(settings.get('Pink Background', 'FFE0E0')),
        'purple': hex_to_rgb(settings.get('Purple Background', 'E0D8F0')),
        'yellow': hex_to_rgb(settings.get('Yellow Background', 'FFF0D0')),
        'milestone': hex_to_rgb(settings.get('Milestone Color', 'F0C040')),
        'critical': hex_to_rgb(settings.get('Critical Text Color', 'CC0000')),
        'near_term': hex_to_rgb(settings.get('Near Term Color', '555555')),
        'mid_term': hex_to_rgb(settings.get('Mid Term Color', 'C8A080')),
        'far_term': hex_to_rgb(settings.get('Far Term Color', 'E8D8C8')),
        'tan_box': hex_to_rgb('C0A080'),
    }
    
    bg_color_map = {
        'pink': colors['pink'],
        'purple': colors['purple'],
        'yellow': colors['yellow'],
    }
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # Background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(0xe0, 0xe0, 0xd0)
    background.line.fill.background()
    
    # Title
    title = settings.get('Title', 'Roadmap')
    add_text_box(slide, Inches(0), Inches(0.15), SLIDE_WIDTH, Inches(0.5),
                 title, font_size=24, bold=True, italic=True, alignment=PP_ALIGN.CENTER)
    
    # Timeline header
    header_top = Inches(0.6)
    header_height = Inches(0.3)
    left_margin = Inches(1.4)
    
    # Add "Fiscal Years" label
    add_rectangle(slide, left_margin, header_top, Inches(1.0), header_height, colors['navy'])
    add_text_box(slide, left_margin, header_top, Inches(1.0), header_height,
                 "Fiscal Years →", font_size=9, bold=True, font_color=RGBColor(255,255,255),
                 alignment=PP_ALIGN.CENTER)
    
    # Build year columns and track positions
    x_pos = left_margin + Inches(1.0)
    year_positions = {}
    
    for _, row in timeline_df.iterrows():
        year = str(row['Year'])
        width = Inches(row['Width (inches)'])
        
        add_rectangle(slide, x_pos, header_top, width, header_height, colors['navy'])
        add_text_box(slide, x_pos, header_top, width, header_height,
                     year, font_size=9, bold=True, font_color=RGBColor(255,255,255),
                     alignment=PP_ALIGN.CENTER)
        
        # Store position for milestone placement
        year_key = year.split("-")[0]  # Use first year for ranges
        year_positions[year_key] = x_pos
        
        x_pos += width
    
    total_grid_width = x_pos - left_margin - Inches(1.0)
    
    # Layout constants
    stc_width = Inches(0.4)
    fta_width = Inches(1.0)
    
    # Process goals and rows
    current_top = Inches(0.95)
    
    for _, goal in goals_df.iterrows():
        goal_id = goal['Goal ID']
        goal_name = goal['Goal Name']
        
        # Add goal header
        add_text_box(slide, Inches(0), current_top, SLIDE_WIDTH, Inches(0.25),
                     goal_name, font_size=12, bold=True, italic=True, alignment=PP_ALIGN.CENTER)
        current_top += Inches(0.3)
        
        # Get rows for this goal
        goal_rows = rows_df[rows_df['Goal ID'] == goal_id].copy()
        
        # Track STC spans
        stc_spans = {}
        for _, row in goal_rows.iterrows():
            stc = row['STC Label']
            if stc not in stc_spans:
                stc_spans[stc] = {'start': current_top, 'count': 0}
            stc_spans[stc]['count'] += 1
        
        # Track current STC for drawing
        current_stc = None
        stc_start_top = current_top
        stc_row_count = 0
        
        for idx, (_, row) in enumerate(goal_rows.iterrows()):
            row_id = row['Row ID']
            stc = row['STC Label']
            fta = row['FTA Label']
            bg_color_name = str(row['Background Color']).lower()
            row_height = Inches(row['Row Height (inches)'])
            
            bg_color = bg_color_map.get(bg_color_name, colors['pink'])
            
            # Handle STC column
            if stc != current_stc:
                # Draw previous STC column if exists
                if current_stc is not None and stc_row_count > 0:
                    stc_height = row_height * stc_row_count
                    add_rectangle(slide, 0, stc_start_top, stc_width, stc_height, colors['navy'])
                    add_text_box(slide, Inches(0.02), stc_start_top + stc_height/2 - Inches(0.1),
                                 stc_width - Inches(0.04), Inches(0.25),
                                 current_stc, font_size=9, bold=True, 
                                 font_color=RGBColor(255,255,255), alignment=PP_ALIGN.CENTER)
                
                current_stc = stc
                stc_start_top = current_top
                stc_row_count = 1
            else:
                stc_row_count += 1
            
            # Draw FTA column
            add_rectangle(slide, stc_width, current_top, fta_width, row_height, colors['navy'])
            add_text_box(slide, stc_width + Inches(0.25), current_top + row_height/2 - Inches(0.1),
                         Inches(0.5), Inches(0.25), fta, font_size=11, bold=True,
                         font_color=RGBColor(255,255,255), alignment=PP_ALIGN.CENTER)
            
            # Draw grid background
            add_rectangle(slide, left_margin, current_top, total_grid_width + Inches(1.0), 
                          row_height, bg_color)
            
            # Draw dashed vertical lines
            for year_key, year_x in year_positions.items():
                line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, year_x, current_top, 
                                               Pt(1), row_height)
                line.fill.solid()
                line.fill.fore_color.rgb = RGBColor(0xcc, 0xcc, 0xcc)
                line.line.fill.background()
            
            # Add milestones for this row
            row_milestones = milestones_df[milestones_df['Row ID'] == row_id]
            
            for _, ms in row_milestones.iterrows():
                year = str(int(ms['Year']))
                v_offset = ms['Vertical Offset']
                text = ms['Text']
                is_critical = str(ms['Is Critical']).lower() == 'yes'
                
                # Find x position
                if year in year_positions:
                    x = year_positions[year] + Inches(0.15)
                elif int(year) >= 2029:
                    x = year_positions.get('2029', list(year_positions.values())[-1]) + Inches(0.3)
                else:
                    continue
                
                y = current_top + Inches(v_offset)
                
                # Add triangle
                add_triangle_milestone(slide, x, y, Inches(0.12), colors['milestone'])
                
                # Add text
                text_color = colors['critical'] if is_critical else RGBColor(0x33, 0x33, 0x33)
                add_text_box(slide, x + Inches(0.14), y, Inches(0.9), Inches(0.15),
                            text, font_size=7, font_color=text_color)
            
            current_top += row_height
        
        # Draw final STC column for this goal
        if current_stc is not None and stc_row_count > 0:
            stc_height = Inches(rows_df[rows_df['Goal ID'] == goal_id]['Row Height (inches)'].iloc[0]) * stc_row_count
            add_rectangle(slide, 0, stc_start_top, stc_width, stc_height, colors['navy'])
            add_text_box(slide, Inches(0.02), stc_start_top + stc_height/2 - Inches(0.1),
                         stc_width - Inches(0.04), Inches(0.25),
                         current_stc, font_size=9, bold=True,
                         font_color=RGBColor(255,255,255), alignment=PP_ALIGN.CENTER)
    
    # Term Arrows
    arrow_top = current_top + Inches(0.1)
    arrow_height = Inches(0.35)
    
    add_rectangle(slide, left_margin, arrow_top, Inches(2.4), arrow_height, colors['near_term'])
    add_text_box(slide, left_margin, arrow_top, Inches(2.4), arrow_height,
                 "Near Term", font_size=11, bold=True, italic=True,
                 font_color=RGBColor(255,255,255), alignment=PP_ALIGN.CENTER)
    
    add_rectangle(slide, left_margin + Inches(2.3), arrow_top, Inches(2.0), arrow_height, colors['mid_term'])
    add_text_box(slide, left_margin + Inches(2.3), arrow_top, Inches(2.0), arrow_height,
                 "Mid Term", font_size=11, bold=True, italic=True,
                 font_color=RGBColor(255,255,255), alignment=PP_ALIGN.CENTER)
    
    add_rectangle(slide, left_margin + Inches(4.2), arrow_top, Inches(3.4), arrow_height, colors['far_term'])
    add_text_box(slide, left_margin + Inches(4.2), arrow_top, Inches(3.4), arrow_height,
                 "Far Term", font_size=11, bold=True, italic=True,
                 font_color=RGBColor(0x33,0x33,0x33), alignment=PP_ALIGN.CENTER)
    
    # Legend
    legend_top = arrow_top + Inches(0.5)
    legend_left = Inches(0.3)
    
    add_text_box(slide, legend_left, legend_top, Inches(1), Inches(0.2),
                 "LEGEND", font_size=8, bold=True, italic=True)
    
    # Strategic Test Capability box
    add_rectangle(slide, legend_left, legend_top + Inches(0.2), Inches(1.3), Inches(0.2), colors['tan_box'])
    add_text_box(slide, legend_left, legend_top + Inches(0.2), Inches(1.3), Inches(0.2),
                 "Strategic Test Capability", font_size=7, font_color=RGBColor(255,255,255),
                 alignment=PP_ALIGN.CENTER)
    
    # Functional Test Area box
    add_rectangle(slide, legend_left, legend_top + Inches(0.45), Inches(1.3), Inches(0.2), colors['navy'])
    add_text_box(slide, legend_left, legend_top + Inches(0.45), Inches(1.3), Inches(0.2),
                 "Functional Test Area", font_size=7, font_color=RGBColor(255,255,255),
                 alignment=PP_ALIGN.CENTER)
    
    # Key Milestone
    add_triangle_milestone(slide, legend_left, legend_top + Inches(0.7), Inches(0.12), colors['milestone'])
    add_text_box(slide, legend_left + Inches(0.15), legend_top + Inches(0.7), Inches(1), Inches(0.15),
                 "Key Milestone", font_size=7)
    
    # Use Cases
    uc_left = legend_left + Inches(1.6)
    add_text_box(slide, uc_left, legend_top + Inches(0.2), Inches(1.5), Inches(0.15),
                 "Use Cases:", font_size=7, bold=True)
    
    uc_texts = ["UC1 – Use case 1 text", "UC2 – Use case 2 text", 
                "UC3 – Use case 3 text", "UC4 – Use case 4 text"]
    for i, uc_text in enumerate(uc_texts):
        add_text_box(slide, uc_left, legend_top + Inches(0.35 + i * 0.15), Inches(1.5), Inches(0.15),
                     uc_text, font_size=7)
    
    # UC color indicators from Excel
    uc_indicator_left = uc_left + Inches(1.7)
    for i, (_, uc) in enumerate(usecases_df.iterrows()):
        uc_id = uc['Use Case ID']
        uc_color = hex_to_rgb(str(uc['Color']))
        
        add_triangle_milestone(slide, uc_indicator_left, legend_top + Inches(0.2 + i * 0.2), 
                               Inches(0.1), uc_color)
        add_text_box(slide, uc_indicator_left + Inches(0.12), legend_top + Inches(0.2 + i * 0.2),
                     Inches(1.2), Inches(0.12), uc_id, font_size=6)
    
    return prs


def main():
    if len(sys.argv) < 2:
        print("Usage: python create_roadmap_from_excel.py <excel_file> [output_file]")
        print("Example: python create_roadmap_from_excel.py roadmap_data.xlsx my_roadmap.pptx")
        sys.exit(1)
    
    excel_path = sys.argv[1]
    output_path = sys.argv[2] if len(sys.argv) > 2 else "roadmap_output.pptx"
    
    print(f"Loading data from {excel_path}...")
    data = load_excel_data(excel_path)
    
    print("Creating roadmap...")
    prs = create_roadmap_from_data(data)
    
    prs.save(output_path)
    print(f"Roadmap saved to {output_path}")


if __name__ == "__main__":
    main()
